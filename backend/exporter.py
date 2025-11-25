from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from bs4 import BeautifulSoup
import os
import re
from datetime import datetime

class DocumentExporter:
    def __init__(self, exports_dir: str = './exports'):
        self.exports_dir = exports_dir
        os.makedirs(exports_dir, exist_ok=True)
    
    def parse_markdown_line(self, line: str, doc: Document):
        """Parse a single line of markdown and add to document"""
        line = line.rstrip()
        
        # Skip empty lines
        if not line.strip():
            return
        
        # Headings
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        # Bullet points
        elif line.strip().startswith('•') or line.strip().startswith('- '):
            text = line.strip()[2:] if line.strip().startswith('- ') else line.strip()[1:]
            p = doc.add_paragraph(text.strip(), style='List Bullet')
            self.apply_inline_formatting(p, text.strip())
        # Numbered lists
        elif re.match(r'^\d+\.\s', line.strip()):
            text = re.sub(r'^\d+\.\s', '', line.strip())
            p = doc.add_paragraph(text, style='List Number')
            self.apply_inline_formatting(p, text)
        # Regular paragraph
        else:
            p = doc.add_paragraph()
            self.apply_inline_formatting(p, line)
    
    def apply_inline_formatting(self, paragraph, text: str):
        """Apply bold, italic, and other inline formatting to paragraph"""
        # Clear existing runs
        paragraph.clear()
        
        # Parse inline markdown
        parts = []
        current_pos = 0
        
        # Find all bold (**text**) and italic (*text* or _text_)
        bold_pattern = r'\*\*(.+?)\*\*'
        italic_pattern = r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)|_(.+?)_'
        code_pattern = r'`(.+?)`'
        
        # Combine patterns
        combined_pattern = f'({bold_pattern})|({italic_pattern})|({code_pattern})'
        
        matches = list(re.finditer(combined_pattern, text))
        
        if not matches:
            # No formatting, just add plain text
            paragraph.add_run(text)
            return
        
        # Process matches
        for match in matches:
            # Add text before match
            if match.start() > current_pos:
                paragraph.add_run(text[current_pos:match.start()])
            
            # Determine match type and add formatted text
            if match.group(1):  # Bold
                run = paragraph.add_run(match.group(2))
                run.bold = True
            elif match.group(3) or match.group(4):  # Italic
                content = match.group(4) if match.group(4) else match.group(5)
                run = paragraph.add_run(content)
                run.italic = True
            elif match.group(6):  # Code
                run = paragraph.add_run(match.group(7))
                run.font.name = 'Courier New'
                run.font.size = Pt(10)
            
            current_pos = match.end()
        
        # Add remaining text
        if current_pos < len(text):
            paragraph.add_run(text[current_pos:])
    
    def html_to_text(self, html_content: str) -> str:
        """Convert HTML to plain text"""
        soup = BeautifulSoup(html_content, 'html.parser')
        return soup.get_text()
    
    def export_docx(self, title: str, content: str, outline: str = None) -> str:
        """Export to .docx file with markdown support"""
        doc = Document()
        
        # Add title
        title_para = doc.add_heading(title, 0)
        
        # Parse content line by line
        lines = content.split('\n')
        for line in lines:
            self.parse_markdown_line(line, doc)
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{title.replace(' ', '_')}_{timestamp}.docx"
        filepath = os.path.join(self.exports_dir, filename)
        
        doc.save(filepath)
        return filepath
    
    def export_pptx(self, title: str, content: str, outline: str = None) -> str:
        """Export to .pptx file with professional formatting"""
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide with professional styling
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Extract main title from content or use provided title
        main_title = title
        lines = content.split('\n')
        for line in lines:
            if line.startswith('# '):
                main_title = line[2:].strip()
                break
        
        title_shape.text = main_title
        subtitle.text = f"Generated with DocForge AI • {datetime.now().strftime('%B %d, %Y')}"
        
        # Style title
        if title_shape.has_text_frame:
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = PptxPt(44)
                paragraph.font.bold = True
        
        # Parse markdown content into structured slides
        sections = []
        current_section = None
        skip_first_h1 = False  # Skip the main title if it's a H1
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Main heading (H1) - Skip if it's the title, otherwise new slide
            if line.startswith('# '):
                if not skip_first_h1:
                    skip_first_h1 = True
                    continue  # Skip the main title
                if current_section and current_section['content']:
                    sections.append(current_section)
                current_section = {
                    'title': line[2:].strip(),
                    'content': [],
                    'type': 'title'  # Section divider slide
                }
            # Subheading (H2) - Content slide with bullets
            elif line.startswith('## '):
                if current_section and current_section['content']:
                    sections.append(current_section)
                current_section = {
                    'title': line[3:].strip(),
                    'content': [],
                    'type': 'content'
                }
            # Sub-subheading (H3) - Treat as bold bullet point
            elif line.startswith('### '):
                if current_section:
                    current_section['content'].append({
                        'text': line[4:].strip(),
                        'level': 0,
                        'bold': True
                    })
            # Bullet points
            elif line.startswith('•') or line.startswith('- '):
                if current_section:
                    text = line[2:] if line.startswith('- ') else line[1:]
                    text = text.strip()
                    # Clean markdown
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
                    text = re.sub(r'\*(.+?)\*', r'\1', text)
                    text = re.sub(r'`(.+?)`', r'\1', text)
                    
                    # Determine indentation level
                    level = 0
                    if line.startswith('  ') or line.startswith('\t'):
                        level = 1
                    
                    current_section['content'].append({
                        'text': text,
                        'level': level,
                        'bold': False
                    })
            # Regular paragraphs - add as bullet if not too long
            elif len(line) > 15 and not line.startswith('#'):
                if current_section:
                    # Clean markdown
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
                    text = re.sub(r'\*(.+?)\*', r'\1', text)
                    text = re.sub(r'`(.+?)`', r'\1', text)
                    
                    # Split long paragraphs into multiple bullets
                    if len(text) > 150:
                        # Split by sentences
                        sentences = text.split('. ')
                        for sentence in sentences:
                            if len(sentence.strip()) > 20:
                                current_section['content'].append({
                                    'text': sentence.strip() + ('.' if not sentence.endswith('.') else ''),
                                    'level': 0,
                                    'bold': False
                                })
                    else:
                        current_section['content'].append({
                            'text': text,
                            'level': 0,
                            'bold': False
                        })
        
        if current_section and current_section['content']:
            sections.append(current_section)
        
        # Create slides from sections
        for section in sections:
            if section['type'] == 'title':
                # Section divider slide (bold title, minimal content)
                title_layout = prs.slide_layouts[5]  # Blank or title only
                slide = prs.slides.add_slide(title_layout)
                
                # Add centered title
                left = PptxInches(1)
                top = PptxInches(3)
                width = PptxInches(8)
                height = PptxInches(1.5)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                tf = title_box.text_frame
                tf.text = section['title']
                
                for paragraph in tf.paragraphs:
                    paragraph.font.size = PptxPt(40)
                    paragraph.font.bold = True
                    paragraph.alignment = 1  # Center
                    
            else:
                # Content slide with bullets
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                title_shape.text = section['title']
                
                # Style title
                if title_shape.has_text_frame:
                    for paragraph in title_shape.text_frame.paragraphs:
                        paragraph.font.size = PptxPt(32)
                        paragraph.font.bold = True
                
                # Add content
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                
                # Add up to 7 points per slide for readability
                for i, item in enumerate(section['content'][:7]):
                    if i == 0:
                        p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = item['text']
                    p.level = item['level']
                    p.font.size = PptxPt(18 if item['level'] == 0 else 16)
                    p.font.bold = item['bold']
                    p.space_before = PptxPt(6)
                    p.space_after = PptxPt(6)
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
        filepath = os.path.join(self.exports_dir, filename)
        
        prs.save(filepath)
        return filepath

exporter = DocumentExporter()
